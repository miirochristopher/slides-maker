from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from dataclasses import dataclass, field
import random
import re
import os

# ------------------------------
# Icon Configuration
# ------------------------------

ICON_DIR = "icons"

ICON_MAP = {
    "title": "computer.png",
    "list": "arrow.png",
    "table": "folder.png",
}

def icon_path(name: str) -> str | None:
    path = os.path.join(ICON_DIR, ICON_MAP.get(name, ""))
    return path if os.path.exists(path) else None


def insert_icon(slide, name, left, top, size=0.4):
    path = icon_path(name)
    if not path:
        return
    slide.shapes.add_picture(
        path,
        Inches(left),
        Inches(top),
        width=Inches(size),
        height=Inches(size),
    )

# ------------------------------
# Branding
# ------------------------------

def random_hex_color(min_val=40, max_val=200):
    return "{:02X}{:02X}{:02X}".format(
        random.randint(min_val, max_val),
        random.randint(min_val, max_val),
        random.randint(min_val, max_val),
    )

@dataclass
class Branding:
    logo_path: str | None = None
    brand_text: str | None = None
    primary_color: str = field(default_factory=random_hex_color)
    accent_color: str = field(default_factory=random_hex_color)
    background_color: str = field(default_factory=lambda: random_hex_color(220, 255))

# ------------------------------
# Models
# ------------------------------

@dataclass
class SlideContent:
    title: str
    lines: list[str]

# ------------------------------
# Parser (Faithful)
# ------------------------------

def parse_faithful_notes(text: str) -> list[SlideContent]:
    slides = []
    current_title = None
    buffer = []

    slide_header_pattern = re.compile(
        r"""^.*?(slide|lecture)\s*\d+(\s*[:—-]\s*.*)?$""",
        re.IGNORECASE,
    )

    for line in text.splitlines():
        line = line.rstrip()
        if not line:
            continue
        if line.lower().startswith("presenter note"):
            continue

        if slide_header_pattern.match(line):
            if current_title:
                slides.append(SlideContent(current_title, buffer))
            current_title = line
            buffer = []
        else:
            buffer.append(line)

    if current_title:
        slides.append(SlideContent(current_title, buffer))

    return slides

# ------------------------------
# Utilities
# ------------------------------

def clean_title(raw_title: str, lecture_title: str) -> str:
    if re.match(r"slide\s*1\s*:", raw_title, re.IGNORECASE):
        return lecture_title

    return re.sub(
        r"^(slide|lecture)\s*\d+\s*[:—-]?\s*",
        "",
        raw_title,
        flags=re.IGNORECASE,
    )

def remove_all_shapes(slide):
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        spTree.remove(shape._element)

def apply_background(slide, hex_color):
    if hex_color.upper() == "000000":
        hex_color = "FFFFFF"
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def should_use_table(lines):
    return len(lines) >= 2 and all(":" in l for l in lines)

# ------------------------------
# Code Detection
# ------------------------------

CODE_PATTERNS = re.compile(
    r"""
    (<html|</|<!DOCTYPE|
     \{|\}|;
     ^def\s|^class\s|import\s|
     console\.log|function\s|
     =>|
     public\s+class|
     System\.out\.println|
     useState\(|useEffect\()
    """,
    re.IGNORECASE | re.VERBOSE,
)

def is_code_block(lines: list[str]) -> bool:
    if not lines:
        return False
    code_like = sum(1 for l in lines if CODE_PATTERNS.search(l))
    return code_like >= max(2, len(lines) // 2)

# ------------------------------
# Renderers
# ------------------------------

def insert_title(slide, title):
    insert_icon(slide, "title", 0.6, 0.55)

    box = slide.shapes.add_textbox(
        Inches(1.1),
        Inches(0.5),
        Inches(7.9),
        Inches(1.2),
    )
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

def insert_text(slide, lines, start_y=1.9):
    insert_icon(slide, "list", 0.6, start_y)

    box = slide.shapes.add_textbox(
        Inches(1.1),
        Inches(start_y),
        Inches(7.7),
        Inches(5.0),
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    font_size = max(16, min(26, int((5 * 72) / max(len(lines), 1)) - 4))

    for i, l in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = l
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0, 0, 0)

def insert_code(slide, lines):
    box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.8),
        Inches(8.4),
        Inches(5.2),
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False

    for i, l in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = l
        p.font.name = "Consolas"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(30, 30, 30)

def insert_table(slide, lines):
    insert_icon(slide, "table", 0.6, 2.2)

    data = []
    for l in lines:
        if ":" in l:
            left, right = l.split(":", 1)
            data.append((left.strip(), right.strip()))

    if not data:
        return

    table = slide.shapes.add_table(
        len(data),
        2,
        Inches(1.1),
        Inches(2.2),
        Inches(7.7),
        Inches(4),
    ).table

    for r, (l, v) in enumerate(data):
        table.cell(r, 0).text = l
        table.cell(r, 1).text = v
        for c in (0, 1):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(22)

# ------------------------------
# Intro Slide
# ------------------------------

def build_intro_slide(prs, branding: Branding):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    remove_all_shapes(slide)
    apply_background(slide, branding.background_color)

    if branding.logo_path:
        slide.shapes.add_picture(
            branding.logo_path,
            Inches(3),
            Inches(2),
            width=Inches(4),
        )

# ------------------------------
# Generator (Faithful, Stable)
# ------------------------------

def generate_pptx_faithful(template, notes, output, branding: Branding):
    prs = Presentation(template)
    slides = parse_faithful_notes(notes)

    lecture_title = branding.brand_text or ""

    while prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    build_intro_slide(prs, branding)

    for slide_data in slides:
        raw_title = slide_data.title.strip()
        title = clean_title(raw_title, lecture_title)

        if re.match(r"slide\s*1\s*:", raw_title, re.IGNORECASE):
            continue

        content_lines = slide_data.lines.copy()

        if not title and content_lines:
            title = content_lines.pop(0)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        remove_all_shapes(slide)
        apply_background(slide, branding.background_color)

        insert_title(slide, title)

        if is_code_block(content_lines):
            insert_code(slide, content_lines)
        elif should_use_table(content_lines):
            insert_table(slide, content_lines)
        else:
            insert_text(slide, content_lines)

    prs.save(output)
